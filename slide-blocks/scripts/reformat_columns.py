from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt
import copy
import sys
import os

def sync_style(source_tf, target_tf, is_title=False):
    target_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    target_tf.word_wrap = True
    target_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    try:
        s_font = source_tf.paragraphs[0].runs[0].font
        s_name, s_size, s_bold, s_italic = s_font.name, s_font.size, s_font.bold, s_font.italic
        s_rgb = s_font.color.rgb if hasattr(s_font.color, 'rgb') else None
    except:
        s_name, s_size, s_bold, s_italic, s_rgb = '微软雅黑', Pt(18) if is_title else Pt(11), is_title, False, None

    for p in target_tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name, r.font.size, r.font.bold, r.font.italic = s_name, s_size, s_bold, s_italic
            if s_rgb: r.font.color.rgb = s_rgb

def expand_to_four(input_path, output_path, new_title="New Item", new_content="Description"):
    prs = Presentation(input_path)
    slide = prs.slides[0]
    for s in [s for s in slide.shapes if s.shape_type == 13 and s.width > 12000000]:
        s.element.getparent().remove(s.element)

    W, G, M, OW = 2650000, 200000, 450000, 3583544
    groups = [[], [], []]
    for s in slide.shapes:
        if "应用领域" in (s.text if hasattr(s, 'text') else ""): continue
        if s.left < 4100000: groups[0].append(s)
        elif s.left < 8000000: groups[1].append(s)
        else: groups[2].append(s)

    src_title_tf = next((s.text_frame for s in slide.shapes if hasattr(s, "text") and "Software 1.0" in s.text), None)
    new_lefts = [M + i*(W+G) for i in range(4)]
    
    for i in range(3):
        g = groups[i]
        if not g: continue
        gl = min([s.left for s in g])
        for s in g:
            s.left, s.width = int(new_lefts[i] + (s.left-gl)*(W/OW)), int(s.width*(W/OW))
            if hasattr(s, "text_frame"): sync_style(src_title_tf or s.text_frame, s.text_frame, "Software" in s.text)

    g3l = min([s.left for s in groups[2]])
    for s in groups[2]:
        ns = slide.shapes._spTree.append(copy.deepcopy(s.element))
        new_shape = slide.shapes[-1]
        new_shape.left = int(new_lefts[3] + (s.left-g3l))
        if hasattr(new_shape, "text_frame"):
            if "Software 3.0" in new_shape.text or "LLM" in new_shape.text: new_shape.text = new_title
            else: new_shape.text = new_content
            sync_style(src_title_tf or new_shape.text_frame, new_shape.text_frame, new_title in new_shape.text)

    prs.save(output_path)
    print(f"Success: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python reformat_columns.py <input.pptx> <output.pptx> [title] [content]")
    else:
        expand_to_four(sys.argv[1], sys.argv[2], sys.argv[3] if len(sys.argv)>3 else "AI Agent", sys.argv[4] if len(sys.argv)>4 else "Autonomous Entity")

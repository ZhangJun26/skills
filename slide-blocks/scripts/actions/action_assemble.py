import os
import copy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def find_title_layout(prs):
    """寻找模板中仅包含标题的版式，或者标准的标题+正文版式"""
    for layout in prs.slide_layouts:
        has_title = False
        has_body = False
        for shape in layout.shapes:
            if not shape.is_placeholder:
                continue
            ph_type = shape.placeholder_format.type
            if ph_type == 1 or ph_type == 3: # TITLE (1) or CENTER_TITLE (3)
                has_title = True
            elif ph_type == 2: # BODY (2)
                has_body = True
        if has_title and not has_body:
            return layout
    return prs.slide_layouts[1] # fallback 到第二个通常是“标题和内容”

def inject_slide(source_slide, target_prs, title_layout):
    """将源幻灯片的内容注入到模板的新幻灯片中"""
    new_slide = target_prs.slides.add_slide(title_layout)
    
    # 1. 注入标题 (使用模板的标题样式)
    source_title_id = -1
    if source_slide.shapes.title:
        source_title_id = source_slide.shapes.title.shape_id
        if new_slide.shapes.title:
            new_slide.shapes.title.text = source_slide.shapes.title.text
            
            # 也可以尝试寻找没有被识别为 title 但是实际是标题的文本框
    
    # 2. 深度克隆正文 Blocks
    for shape in source_slide.shapes:
        # 跳过源标题（因为已经注入了）
        if shape.shape_id == source_title_id:
            continue
            
        # 智能过滤：跳过全屏背景图 (Width > 12000000 EMU)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.width > 12000000:
            continue
            
        # 核心引擎：XML 级深度克隆
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.append(new_el)

def assemble():
    base_dir = r"D:\AI\.claude\素材"
    template_path = os.path.join(base_dir, "浅色底模板.pptx")
    hami_path = os.path.join(base_dir, "完整版-哈密市医疗健康数智融合发展规划建议.260202.pptx")
    ai_path = os.path.join(base_dir, "完整版-行业会议-医疗AI的思考与实践-浅色底.251031.pptx")
    
    target_prs = Presentation(template_path)
    # 删除模板自带的示例页（如果需要纯净的开始）
    # xml_slides = target_prs.slides._sldIdLst
    # for slide_id in list(xml_slides): xml_slides.remove(slide_id)
    
    title_layout = find_title_layout(target_prs)
    
    # --- 组装块 1：哈密公司介绍 (Slide 2-6, index 1-5) ---
    print("正在注入: 模块 1 [哈密公司介绍]")
    hami_prs = Presentation(hami_path)
    for i in range(1, 6):
        inject_slide(hami_prs.slides[i], target_prs, title_layout)
        print(f"  - 成功克隆第 {i+1} 页")

    # --- 组装块 2：医疗 AI 场景 (Slide 19-25, index 18-24) ---
    print("正在注入: 模块 2 [医疗AI场景]")
    ai_prs = Presentation(ai_path)
    for i in range(18, 25):
        inject_slide(ai_prs.slides[i], target_prs, title_layout)
        print(f"  - 成功克隆第 {i+1} 页")

    # 保存组装成品
    output_path = r"D:\AI\.claude\已组装-公司介绍与AI场景.pptx"
    target_prs.save(output_path)
    print(f"\n拼接完成！文件已保存至: {output_path}")

if __name__ == "__main__":
    assemble()

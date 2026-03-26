# -*- coding: utf-8 -*-
"""
assemble_template.py — 模板优先的 PPT 组装脚本 v2

策略：
  - 纯模板页（封面/过渡页/封底）：直接复制模板指定页，ExecuteMso 保留源格式粘贴
  - 内容页：先复制模板 P3 建立背景+标题栏，再从源文件复制内容形状，
    切换到 ppViewSlide（幻灯片编辑区）后 ExecuteMso 保留源格式粘贴 → 颜色不变

Plan 项格式：
  纯模板页：{"template_page": N, "replace_title": "xxx"}
  内容页：  {"src": "路径", "page": N, "replace_title": "xxx"}
"""

import time
import sys
import subprocess
from pathlib import Path
from datetime import datetime

# ─── PowerPoint 启动路径（WPS 劫持 COM 时的绕过方案） ─────────────────────────
# WPS 会把自己注册为 PowerPoint.Application 的 COM 提供方，导致 Dispatch() 拿到 WPS。
# 绕过方式：直接启动 POWERPNT.EXE，再用 GetActiveObject 附加真正的 PowerPoint 实例。
# 如果机器上只有 WPS，找不到任何路径，系统会回退到 Dispatch("PowerPoint.Application")。
_POWERPNT_CANDIDATES = [
    r"C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE",
    r"C:\Program Files (x86)\Microsoft Office\root\Office16\POWERPNT.EXE",
    r"C:\Program Files\Microsoft Office\Office16\POWERPNT.EXE",
    r"C:\Program Files (x86)\Microsoft Office\Office16\POWERPNT.EXE",
    r"C:\Program Files\Microsoft Office\root\Office15\POWERPNT.EXE",
    r"C:\Program Files (x86)\Microsoft Office\root\Office15\POWERPNT.EXE",
    r"C:\Program Files\Microsoft Office\Office15\POWERPNT.EXE",
    r"C:\Program Files (x86)\Microsoft Office\Office15\POWERPNT.EXE",
]
POWERPNT_EXE = next((p for p in _POWERPNT_CANDIDATES if Path(p).exists()), None)

# ─── 路径配置 ─────────────────────────────────────────────────────────────────

_BASE_DIR = Path(__file__).parent.parent  # engine/ 的上一层 = SlideBlocks 根目录

TEMPLATE_PATH = _BASE_DIR / "模板" / "科技风（深色底）.pptx"
TEMPLATE_CONTENT_PAGE          = 3  # 带标题栏内容页（模板P3）
TEMPLATE_CONTENT_PAGE_NO_TITLE = 4  # 无标题栏内容页（模板P4）

OUTPUT_DIR = _BASE_DIR / "输出"

# 标题区域阈值（pt）：top 小于此值的形状视为标题栏，不复制进输出
# win32com shape.Top 单位是 pt（不是 EMU）
TITLE_THRESHOLD_PT = 65

# PowerPoint ViewType 常量
PP_VIEW_NORMAL = 1   # 普通视图（缩略图面板+编辑区）
PP_VIEW_SLIDE  = 9   # 幻灯片视图（全屏编辑单张幻灯片，编辑区获得焦点）

# ─── 辅助函数 ─────────────────────────────────────────────────────────────────

def get_source_title(src_slide):
    """从源幻灯片提取标题文字，同时返回该形状的索引（用于精准排除）。
    返回：(title_text, shape_index) 或 (None, None)
    """
    # 优先：标题占位符（ppPlaceholderTitle = 1）
    for j in range(1, src_slide.Shapes.Count + 1):
        shape = src_slide.Shapes(j)
        try:
            if shape.PlaceholderFormat.Type == 1 and shape.HasTextFrame:
                t = shape.TextFrame.TextRange.Text.strip()
                if t:
                    return t, j
        except Exception:
            pass
    # 备选：标题区域（top < TITLE_THRESHOLD_PT）内有文字、top 最小的形状
    candidates = []
    for j in range(1, src_slide.Shapes.Count + 1):
        shape = src_slide.Shapes(j)
        if shape.Top < TITLE_THRESHOLD_PT:
            try:
                if shape.HasTextFrame:
                    t = shape.TextFrame.TextRange.Text.strip()
                    if t:
                        candidates.append((shape.Top, j, t))
            except Exception:
                pass
    if candidates:
        candidates.sort()
        _, j, t = candidates[0]
        return t, j
    return None, None


def get_content_indices(src_slide, exclude_idx=None):
    """返回要复制的形状索引列表。
    排除规则：
    - 标题占位符（ppPlaceholderTitle=1）
    - exclude_idx 指定的标题文字形状
    - 标题区（top < TITLE_THRESHOLD_PT）内的图片/图形形状（type=13，msoPicture），
      即 logo、图标等装饰图片；文字形状保留（如标题区右侧的矩形文字框）
    """
    MSO_PICTURE = 13  # msoPicture
    indices = []
    for j in range(1, src_slide.Shapes.Count + 1):
        shape = src_slide.Shapes(j)
        # 跳过标题占位符
        try:
            if shape.PlaceholderFormat.Type == 1:
                continue
        except Exception:
            pass
        # 跳过检测到的标题形状
        if j == exclude_idx:
            continue
        # 标题区内的图片形状（logo/图标）一律排除
        if shape.Top < TITLE_THRESHOLD_PT and shape.Type == MSO_PICTURE:
            continue
        indices.append(j)
    return indices


def set_template_title(target_slide, text):
    """将文字写入模板标题文本框（top 最小、且当前有非空文字的形状）。
    必须过滤掉空文字的背景矩形，否则会写错目标。
    """
    best, best_top = None, float("inf")
    for j in range(1, target_slide.Shapes.Count + 1):
        shape = target_slide.Shapes(j)
        try:
            if shape.HasTextFrame:
                existing = shape.TextFrame.TextRange.Text.strip()
                if existing and shape.Top < best_top:  # 必须有非空文字
                    best_top = shape.Top
                    best = shape
        except Exception:
            pass
    if best:
        try:
            best.TextFrame.TextRange.Text = text
        except Exception as e:
            print(f"    [警告] 写入标题失败: {e}")


def paste_slide_with_source_format(pptApp, pres, count_before):
    """在幻灯片缩略图面板（PP_VIEW_NORMAL）上下文里，ExecuteMso 粘贴整页（保留源格式）。"""
    pres.Windows(1).Activate()
    pres.Windows(1).ViewType = PP_VIEW_NORMAL
    time.sleep(0.3)
    if pres.Slides.Count > 0:
        pres.Windows(1).View.GotoSlide(pres.Slides.Count)
    pptApp.CommandBars.ExecuteMso("PasteSourceFormatting")
    deadline = time.time() + 5.0
    while pres.Slides.Count == count_before and time.time() < deadline:
        time.sleep(0.1)
    if pres.Slides.Count == count_before:
        print(f"    [警告] ExecuteMso 超时，降级普通粘贴")
        pres.Slides.Paste(pres.Slides.Count + 1)


def paste_shapes_with_source_format(pptApp, pres, current_idx):
    """在幻灯片编辑区（PP_VIEW_SLIDE）上下文里，ExecuteMso 粘贴形状（保留源格式、颜色不变）。"""
    pres.Windows(1).Activate()
    pres.Windows(1).View.GotoSlide(current_idx)
    pres.Windows(1).ViewType = PP_VIEW_SLIDE   # 切到全屏编辑视图，编辑区获得焦点
    time.sleep(0.3)
    pptApp.CommandBars.ExecuteMso("PasteSourceFormatting")
    time.sleep(0.4)
    pres.Windows(1).ViewType = PP_VIEW_NORMAL  # 切回普通视图
    time.sleep(0.2)


# ─── 颜色修复（深色底来源 → 浅色底模板适配） ─────────────────────────────────────

_DARK_HEX = '262626'              # 深灰，替换浅色实色文字
# 渐变文字替换色：[起始色, 结束色]，中间停止点线性插值
# RGB(1,101,255) → RGB(0,184,243)：蓝到青蓝
_GRADIENT_FALLBACK_COLORS = ['0165FF', '00B8F3']

# 浅色预设颜色名（prstClr val="white" 等），全部替换为深色
_LIGHT_PRESET_COLORS = {
    'white', 'ltGray', 'silver', 'gainsboro', 'ghostWhite',
    'snow', 'ivory', 'floralWhite', 'honeydew', 'azure', 'aliceBlue',
    'lavenderBlush', 'mistyRose', 'seashell', 'linen', 'oldLace',
    'mintCream', 'antiqueWhite', 'cornsilk', 'lemonChiffon',
}

def _is_light_hex(hex_val):
    """判断 6 位十六进制颜色是否为浅色（RGB 各分量 >= 200）"""
    try:
        r, g, b = int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
        return r >= 200 and g >= 200 and b >= 200
    except Exception:
        return False


def _hex_luminance(hex_val):
    """计算 6 位十六进制颜色的感知亮度（0-255）"""
    r = int(hex_val[0:2], 16)
    g = int(hex_val[2:4], 16)
    b = int(hex_val[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def _srgbclr_effective_luminance(srgbclr_el):
    """
    计算 <a:srgbClr> 元素的实际感知亮度（0-255），考虑 lumMod/lumOff 子元素。
    例：val="5B9BD5" + lumOff=95000 → L 接近 100% → 近白色。
    无法计算时返回 None（调用方回退到 _is_light_hex）。
    """
    import colorsys
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    val = srgbclr_el.get('val', '')
    if len(val) != 6:
        return None
    try:
        r = int(val[0:2], 16) / 255
        g = int(val[2:4], 16) / 255
        b = int(val[4:6], 16) / 255
        h, l, s = colorsys.rgb_to_hls(r, g, b)

        lum_mod = srgbclr_el.find(f'{{{NS_A}}}lumMod')
        lum_off = srgbclr_el.find(f'{{{NS_A}}}lumOff')
        if lum_mod is not None:
            l = l * (int(lum_mod.get('val', '100000')) / 100000)
        if lum_off is not None:
            l = l + (int(lum_off.get('val', '0')) / 100000)
        l = max(0.0, min(1.0, l))

        r2, g2, b2 = colorsys.hls_to_rgb(h, l, s)
        return 0.299 * r2 * 255 + 0.587 * g2 * 255 + 0.114 * b2 * 255
    except Exception:
        return None


def _shape_has_dark_fill(shape, threshold=140):
    """
    判断形状是否有深色背景填充（深色 → 白字应保留）。
    threshold=140：亮度低于此值视为深色（0-255 范围）。
    支持 SOLID（纯色）和 GRADIENT（渐变，取所有色标平均亮度）。
    返回 False 时表示浅色/透明/不确定 → 白字应改深色。
    """
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
        fill = shape.fill
        ft = fill.type
        if ft is None or ft == MSO_FILL.BACKGROUND:
            return False  # 透明，底色是模板浅色背景

        if ft == MSO_FILL.SOLID:
            fc = fill.fore_color
            if fc.type == MSO_COLOR_TYPE.RGB:
                return _hex_luminance(str(fc.rgb)) < threshold

        if ft == MSO_FILL.GRADIENT:
            # python-pptx 渐变 API 不完整，直接读 XML 色标
            NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            grad = shape._element.find(f'.//{{{NS_A}}}gradFill')
            if grad is not None:
                stops = grad.findall(f'.//{{{NS_A}}}gs')
                lums = []
                alphas = []
                for gs in stops:
                    clr = gs.find(f'{{{NS_A}}}srgbClr')
                    if clr is not None:
                        val = clr.get('val', '')
                        if len(val) == 6:
                            lums.append(_hex_luminance(val))
                        # 检测 alpha 透明度
                        alpha_el = clr.find(f'{{{NS_A}}}alpha')
                        if alpha_el is not None:
                            alphas.append(int(alpha_el.get('val', '100000')) / 100000)
                # 渐变平均 alpha < 50% → 实际接近透明，视为非深色背景
                if alphas and (sum(alphas) / len(alphas)) < 0.5:
                    return False
                if lums:
                    return (sum(lums) / len(lums)) < threshold

        # 图案/主题色 → 保守处理（返回 False）
        return False
    except Exception:
        return False


# 深色底模板里用作"浅色文字"的主题色槽位，粘到浅色底模板后会变成白色（不可见）
_LIGHT_SCHEME_COLORS = {'bg1', 'lt1', 'bg2', 'lt2'}


def _fix_solidfill_el(sf_el, dark_hex):
    """修复一个 <a:solidFill> 元素：浅色/空 → 深色"""
    from lxml import etree
    NS_A      = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    PRSTCLR   = f'{{{NS_A}}}prstClr'
    SRGBCLR   = f'{{{NS_A}}}srgbClr'
    SCHEMECLR = f'{{{NS_A}}}schemeClr'

    children = list(sf_el)
    if not children:
        # 空的 solidFill → 设为深色
        etree.SubElement(sf_el, SRGBCLR).set('val', dark_hex)
        return
    child = children[0]
    if child.tag == PRSTCLR and child.get('val', '') in _LIGHT_PRESET_COLORS:
        sf_el.remove(child)
        etree.SubElement(sf_el, SRGBCLR).set('val', dark_hex)
    elif child.tag == SRGBCLR:
        # 先计算含 lumMod/lumOff 的实际亮度，再回退到基础色判断
        eff_lum = _srgbclr_effective_luminance(child)
        is_light = (eff_lum is not None and eff_lum > 200) or \
                   (eff_lum is None and _is_light_hex(child.get('val', '')))
        if is_light:
            child.set('val', dark_hex)
            # 移除 lumMod/lumOff（已替换为纯色，修饰子元素无意义）
            for mod_tag in [f'{{{NS_A}}}lumMod', f'{{{NS_A}}}lumOff', f'{{{NS_A}}}lumClamp']:
                mod_el = child.find(mod_tag)
                if mod_el is not None:
                    child.remove(mod_el)
    elif child.tag == SCHEMECLR:
        is_light_slot = child.get('val', '') in _LIGHT_SCHEME_COLORS
        # 任意主题色 + lumOff >= 40000（向白色推移 40%+）→ 实际接近白色
        lum_off_el = child.find(f'{{{NS_A}}}lumOff')
        has_high_lumoff = lum_off_el is not None and int(lum_off_el.get('val', '0')) >= 40000
        if is_light_slot or has_high_lumoff:
            sf_el.remove(child)
            etree.SubElement(sf_el, SRGBCLR).set('val', dark_hex)


def _fix_gradfill_el(gf_el, dark_hex):
    """
    修复文字 <a:gradFill>：若所有色标均为浅色，整个渐变替换为深色 solidFill。
    有深色色标（刻意设计的彩色渐变）则保留不改。
    """
    from lxml import etree
    NS_A    = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SRGBCLR = f'{{{NS_A}}}srgbClr'
    SCHEMECLR = f'{{{NS_A}}}schemeClr'
    parent  = gf_el.getparent()
    if parent is None:
        return

    # 收集所有色标的有效亮度
    lums = []
    for gs in gf_el.iter(f'{{{NS_A}}}gs'):
        clr = gs.find(SRGBCLR)
        if clr is not None:
            val = clr.get('val', '')
            if len(val) == 6:
                # 考虑 lumMod/lumOff 的实际亮度
                eff = _srgbclr_effective_luminance(clr)
                lums.append(eff if eff is not None else _hex_luminance(val))
        sc = gs.find(SCHEMECLR)
        if sc is not None:
            is_light_slot = sc.get('val', '') in _LIGHT_SCHEME_COLORS
            lum_off_el = sc.find(f'{{{NS_A}}}lumOff')
            has_high_lumoff = lum_off_el is not None and int(lum_off_el.get('val', '0')) >= 40000
            if is_light_slot or has_high_lumoff:
                lums.append(255)  # 视为浅色

    # 全部色标都是浅色（平均亮度 > 200）→ 替换为渐变色，保留原有渐变结构（角度/方向不变）
    if lums and (sum(lums) / len(lums)) > 200:
        stops = [gs for gs in gf_el.iter(f'{{{NS_A}}}gs')]
        n = len(stops)
        c1 = _GRADIENT_FALLBACK_COLORS[0]
        c2 = _GRADIENT_FALLBACK_COLORS[-1]
        for i, gs in enumerate(stops):
            # 线性插值：0 → c1，n-1 → c2
            t = i / max(n - 1, 1)
            r = int(int(c1[0:2], 16) * (1 - t) + int(c2[0:2], 16) * t)
            g = int(int(c1[2:4], 16) * (1 - t) + int(c2[2:4], 16) * t)
            b = int(int(c1[4:6], 16) * (1 - t) + int(c2[4:6], 16) * t)
            interp = f'{r:02X}{g:02X}{b:02X}'
            # 清空色标内原有颜色子元素，换成新 srgbClr
            for old_clr in list(gs):
                gs.remove(old_clr)
            etree.SubElement(gs, SRGBCLR).set('val', interp)


def _fix_text_colors_xml(slide_element, dark_hex=_DARK_HEX):
    """
    直接操作 XML，修复文本 rPr / defRPr 里的浅色颜色。
    处理：srgbClr 近白色、prstClr white、schemeClr bg1/lt1、gradFill 全浅色。
    只改文字属性（rPr/defRPr），不动形状背景填充。
    """
    NS_A      = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SOLIDFILL = f'{{{NS_A}}}solidFill'
    GRADFILL  = f'{{{NS_A}}}gradFill'
    RPR       = f'{{{NS_A}}}rPr'
    DEFRPR    = f'{{{NS_A}}}defRPr'

    for tag in (RPR, DEFRPR):
        for el in slide_element.iter(tag):
            for sf in el.findall(SOLIDFILL):
                try:
                    _fix_solidfill_el(sf, dark_hex)
                except Exception:
                    pass
            for gf in el.findall(GRADFILL):
                try:
                    _fix_gradfill_el(gf, dark_hex)
                except Exception:
                    pass


# 模块级变量：当前处理 slide 的深色形状区域（left, top, right, bottom），单位 EMU
# 由 fix_colors_for_light_template 在处理每页前更新
_dark_regions: list = []


def _shape_on_dark_region(shape):
    """
    判断透明/无填充形状是否叠在深色兄弟形状上（slide 级位置检测）。
    用于解决 TEXT_BOX 透明填充但视觉上位于深色面板上的误改问题（P49 类型）。
    """
    if not _dark_regions:
        return False
    try:
        sl = shape.left;  st = shape.top
        sr = sl + shape.width;  sb = st + shape.height
        for (rl, rt, rr, rb) in _dark_regions:
            # 重叠面积超过形状面积 30% → 视为在深色区域上
            ox = max(0, min(sr, rr) - max(sl, rl))
            oy = max(0, min(sb, rb) - max(st, rt))
            overlap = ox * oy
            area = max((sr - sl) * (sb - st), 1)
            if overlap / area > 0.65:
                return True
    except Exception:
        pass
    return False


def _fix_gradient_text_xml(element):
    """
    只修复元素内的渐变文字（gradFill）→ 标准蓝。
    用于深色背景形状：保留实色白字，但渐变近白色文字必须改（lumOff 设计的近白色在浅色底无效）。
    """
    NS_A     = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    GRADFILL = f'{{{NS_A}}}gradFill'
    RPR      = f'{{{NS_A}}}rPr'
    DEFRPR   = f'{{{NS_A}}}defRPr'
    for tag in (RPR, DEFRPR):
        for el in element.iter(tag):
            for gf in el.findall(GRADFILL):
                try:
                    _fix_gradfill_el(gf, _DARK_HEX)
                except Exception:
                    pass


def _fix_shape_text_colors_smart(shape, dark_hex=_DARK_HEX, _sibling_dark=None):
    """
    修复单个形状的文字颜色（深色底来源 → 浅色底）：
    - 形状有显式深色填充 → 保留白字（只修渐变文字）
    - 形状透明 + 与当前 GROUP 内的兄弟深色形状重叠 ≥ 65% → 视觉上在深色背景上，保白
    - 其他所有情况 → 全部改为深色

    关键设计：_sibling_dark 只包含同一 GROUP 层级的直接兄弟深色形状。
    这避免了全局 _dark_regions 的跨组误判问题：
    - 架构图 GROUP：深色格子在子 GROUP 里，外层兄弟无深色形状 → 透明标题改黑
    - Agent Framework GROUP：深色面板和透明标签是同层兄弟，直接叠放 → 标签保白
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # 构建本 GROUP 直接子级的深色形状区域（不递归，只看直接兄弟）
            sibling_dark = []
            for child in shape.shapes:
                if _shape_has_dark_fill(child):
                    sibling_dark.append((child.left, child.top,
                                         child.left + child.width,
                                         child.top + child.height))
            for child in shape.shapes:
                _fix_shape_text_colors_smart(child, dark_hex, _sibling_dark=sibling_dark)
            return
    except Exception:
        pass

    if _shape_has_dark_fill(shape):
        # 深色填充 → 背景深色，白字正确，只修渐变文字
        try:
            _fix_gradient_text_xml(shape._element)
        except Exception:
            pass
        return

    # 透明形状：检查是否与兄弟深色形状大面积重叠
    if _sibling_dark:
        try:
            from pptx.enum.dml import MSO_FILL
            ft = shape.fill.type
            is_transparent = ft is None or ft == MSO_FILL.BACKGROUND
        except Exception:
            is_transparent = True
        if is_transparent:
            try:
                sl = shape.left; st = shape.top
                sr = sl + shape.width; sb = st + shape.height
                area = max((sr - sl) * (sb - st), 1)
                for (rl, rt, rr, rb) in _sibling_dark:
                    ox = max(0, min(sr, rr) - max(sl, rl))
                    oy = max(0, min(sb, rb) - max(st, rt))
                    if (ox * oy) / area > 0.65:
                        # 叠在深色兄弟形状上，保留白字（只修渐变）
                        try:
                            _fix_gradient_text_xml(shape._element)
                        except Exception:
                            pass
                        return
            except Exception:
                pass

    try:
        _fix_text_colors_xml(shape._element, dark_hex)
    except Exception:
        pass


def _fix_shape_fills(shape):
    """修复形状背景填充：显式浅色实心填充 → 透明（递归处理组合形状）"""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        st = shape.shape_type
        if st in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
            return
        if st == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                _fix_shape_fills(child)
            return
    except Exception:
        pass
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
        fill = shape.fill
        if fill.type == MSO_FILL.SOLID:
            fc = fill.fore_color
            if fc.type == MSO_COLOR_TYPE.RGB and _is_light_hex(str(fc.rgb)):
                fill.background()
    except Exception:
        pass


# 大面积深色形状透明化：面积 > 此比例 → 加透明度（避免深色底来源的色块在浅色底上过于突兀）
_LARGE_DARK_AREA_THRESHOLD = 0.008  # 占幻灯片面积 0.8%（覆盖小圆盘等设计元素）
_LARGE_DARK_ALPHA = 15000           # 15% 不透明（85% 透明），深蓝 → 很淡的色调

# 大面积浅色形状透明化：浅色底来源 → 深色底模板，浅色大卡片在深色背景上太亮
_LARGE_LIGHT_AREA_THRESHOLD = 0.05  # 占幻灯片面积 5%（针对大卡片/色块）
_LARGE_LIGHT_ALPHA = 20000          # 20% 不透明（80% 透明），白色 → 半透明玻璃感


def _lighten_large_dark_shapes(slide, slide_area):
    """
    对面积超过阈值的大面积深色实色形状，添加透明度。
    深色底的大色块（如架构图底座）在浅色底上会显得很突兀，加透明度后变成淡色调。
    只处理实色填充（SOLID）；渐变形状暂不处理。
    """
    from lxml import etree
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    def _add_alpha(shape):
        try:
            from pptx.enum.dml import MSO_FILL
            if shape.fill.type != MSO_FILL.SOLID:
                return
            # 找形状填充的 solidFill 元素（在 spPr 下，不是文字 rPr 下）
            sp_el = shape._element
            # 只找 spPr 直接子孙的 solidFill，避免误改文字颜色
            sp_pr = sp_el.find(f'.//{{{NS_A}}}spPr') or sp_el  # 兼容不同命名
            # 直接搜 p:spPr 下的 a:solidFill
            NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            for spPr in sp_el.iter(f'{{{NS_P}}}spPr', f'{{{NS_A}}}spPr'):
                sf = spPr.find(f'{{{NS_A}}}solidFill')
                if sf is None:
                    continue
                children = list(sf)
                if not children:
                    continue
                clr = children[0]
                # 移除已有 alpha，插入新 alpha
                old = clr.find(f'{{{NS_A}}}alpha')
                if old is not None:
                    clr.remove(old)
                a_el = etree.SubElement(clr, f'{{{NS_A}}}alpha')
                a_el.set('val', str(_LARGE_DARK_ALPHA))
                return  # 只改第一个 spPr 里的填充
        except Exception:
            pass

    def _process_shapes(shapes):
        for shape in shapes:
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    _process_shapes(shape.shapes)  # 递归进 GROUP
                    continue
                area_pct = (shape.width * shape.height) / max(slide_area, 1)
                if area_pct < _LARGE_DARK_AREA_THRESHOLD:
                    continue
                if not _shape_has_dark_fill(shape):
                    continue
                # 有文字内容的深色形状是内容元素（标签/徽章等），不加透明度。
                # 加透明度后白字会印在近白背景上变得不可见。
                try:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        continue
                except Exception:
                    pass
                _add_alpha(shape)
            except Exception:
                pass

    _process_shapes(slide.shapes)


def _make_large_light_shapes_transparent(slide, slide_area):
    """
    浅色底来源 → 深色底模板：对大面积无显式填充的形状注入半透明白色填充。

    这类形状（如大卡片、背景色块）fill 来自主题 style，没有 spPr 填充声明，
    视觉上呈现浅色/白色。在深色底模板上显得突兀，加高透明度后呈现玻璃感。
    """
    from lxml import etree
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    _FILL_TAGS = {f'{{{NS_A}}}{t}' for t in
                  ('solidFill', 'gradFill', 'blipFill', 'pattFill', 'noFill', 'grpSpFill')}

    def _has_explicit_fill(sp_el):
        """检查 spPr 是否有显式填充声明（有则不动）。"""
        spPr = sp_el.find(f'{{{NS_A}}}spPr')
        if spPr is None:
            return False
        return any(child.tag in _FILL_TAGS for child in spPr)

    def _inject_transparent_fill(sp_el):
        """向 spPr 注入带透明度的白色实色填充。"""
        spPr = sp_el.find(f'{{{NS_A}}}spPr')
        if spPr is None:
            spPr = etree.SubElement(sp_el, f'{{{NS_A}}}spPr')
        sf   = etree.SubElement(spPr, f'{{{NS_A}}}solidFill')
        srgb = etree.SubElement(sf,   f'{{{NS_A}}}srgbClr')
        srgb.set('val', 'FFFFFF')
        alpha = etree.SubElement(srgb, f'{{{NS_A}}}alpha')
        alpha.set('val', str(_LARGE_LIGHT_ALPHA))

    def _process(shapes):
        for shape in shapes:
            try:
                if hasattr(shape, 'shapes'):   # GROUP
                    _process(shape.shapes)
                    continue
                area_pct = (shape.width * shape.height) / max(slide_area, 1)
                if area_pct < _LARGE_LIGHT_AREA_THRESHOLD:
                    continue
                if _has_explicit_fill(shape._element):
                    continue
                _inject_transparent_fill(shape._element)
            except Exception:
                pass

    _process(slide.shapes)


def fix_colors_for_light_template(pptx_path, slide_indices):
    """
    对指定页（1-based）修复深色来源的颜色：
      - 文字颜色（XML 级）：prstClr white / srgbClr 近白色 → 深灰 #262626
      - 形状填充（API 级）：显式浅色实心填充 → 透明
    在 win32com 保存完成后、文件关闭后调用。
    """
    if not slide_indices:
        return
    try:
        from pptx import Presentation
    except ImportError:
        print("  [警告] 未安装 python-pptx，跳过颜色修复（pip install python-pptx）")
        return

    prs = Presentation(str(pptx_path))
    for idx in slide_indices:
        slide = prs.slides[idx - 1]
        for shape in slide.shapes:             # 文字修复：深色填充保白，其余全改黑
            _fix_shape_text_colors_smart(shape)
        for shape in slide.shapes:             # API 级填充修复：浅色实心填充 → 透明
            _fix_shape_fills(shape)
        slide_area = prs.slide_width * prs.slide_height
        _lighten_large_dark_shapes(slide, slide_area)  # 大面积深色形状 → 加透明度
    prs.save(str(pptx_path))
    print(f"  [颜色修复] 已修复第 {slide_indices} 页（深色底来源 → 浅色底适配）")


def _fix_text_to_white_xml(element):
    """
    将元素内非白色文字改为白色（XML 级，用于浅色底来源 → 深色底模板适配）。
    只处理 rPr / defRPr 中的 solidFill。
    规则：透明/无填充背景下，任何非近白色文字 → 白色，确保深色底可见。
    """
    from lxml import etree
    NS_A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SF     = f'{{{NS_A}}}solidFill'
    SRGB   = f'{{{NS_A}}}srgbClr'
    PRST   = f'{{{NS_A}}}prstClr'
    RPR    = f'{{{NS_A}}}rPr'
    DEFRPR = f'{{{NS_A}}}defRPr'
    # 近白色预设：这些颜色在深色底上本已可见，保留
    _NEAR_WHITE_PRESET = {'white', 'ltGray', 'lightGray', 'silver', 'gainsboro'}

    for tag in (RPR, DEFRPR):
        for el in element.iter(tag):
            for sf in el.findall(SF):
                children = list(sf)
                if not children:
                    # 空 solidFill（继承默认色）→ 深色底下继承色为黑，改白
                    etree.SubElement(sf, SRGB).set('val', 'FFFFFF')
                    continue
                child = children[0]
                if child.tag == PRST and child.get('val', '') not in _NEAR_WHITE_PRESET:
                    # 非近白预设（黑、蓝、红等）→ 改白
                    sf.remove(child)
                    etree.SubElement(sf, SRGB).set('val', 'FFFFFF')
                elif child.tag == SRGB:
                    val = child.get('val', '')
                    # 非近白 srgbClr（任一通道 < 200）→ 改白
                    if len(val) == 6 and not _is_light_hex(val):
                        child.set('val', 'FFFFFF')


def _fix_shape_text_to_dark_template(shape):
    """
    智能处理单个形状（浅色底来源 → 深色底模板方向）：
    - 透明/无填充形状：深色文字 → 白色（背景将变为深色模板底）
    - 有显式填充形状：保持原样（形状自身提供视觉背景）
    递归处理组合形状。
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                _fix_shape_text_to_dark_template(child)
            return
    except Exception:
        pass

    try:
        from pptx.enum.dml import MSO_FILL
        ft = shape.fill.type
        is_transparent = (ft is None or ft == MSO_FILL.BACKGROUND)
    except Exception:
        is_transparent = True

    if not is_transparent:
        return  # 有填充 → 保持原文字颜色

    try:
        _fix_text_to_white_xml(shape._element)
    except Exception:
        pass


def fix_colors_for_dark_template(pptx_path, slide_indices):
    """
    对指定页（1-based）修复浅色来源的颜色（适配深色底模板）：
      - 透明/无填充形状中的深色文字 → 白色
      - 有显式填充形状不改（形状自身提供背景色和对比度）
    在 win32com 保存完成后、文件关闭后调用。
    """
    if not slide_indices:
        return
    try:
        from pptx import Presentation
    except ImportError:
        print("  [警告] 未安装 python-pptx，跳过颜色修复（pip install python-pptx）")
        return

    prs = Presentation(str(pptx_path))
    for idx in slide_indices:
        slide = prs.slides[idx - 1]
        slide_area = slide.shapes._spTree.getparent().getparent()
        try:
            from pptx.util import Emu
            w = prs.slide_width
            h = prs.slide_height
            slide_area_val = w * h
        except Exception:
            slide_area_val = 1
        # Step 1：大面积浅色形状 → 半透明（先于文字修复，避免干扰后续判断）
        _make_large_light_shapes_transparent(slide, slide_area_val)
        # Step 2：透明/无填充形状中的非白色文字 → 白色
        for shape in slide.shapes:
            _fix_shape_text_to_dark_template(shape)
    prs.save(str(pptx_path))
    print(f"  [颜色修复] 已修复第 {slide_indices} 页（浅色底来源 → 深色底适配）")


# ─── 主流程 ───────────────────────────────────────────────────────────────────

def _get_ppt_app():
    """获取真正的 PowerPoint COM 实例（绕过 WPS COM 劫持）。"""
    import win32com.client

    exe = Path(POWERPNT_EXE) if POWERPNT_EXE else None

    if exe and exe.exists():
        # 直接启动真正的 POWERPNT.EXE，再附加
        subprocess.Popen([str(exe)])
        # WPS 可能已在 ROT 中注册为 PowerPoint.Application，
        # 需等待真正的 PowerPoint 注册后覆盖。用路径判断而非版本号。
        for _ in range(40):          # 最多等 20 秒
            time.sleep(0.5)
            try:
                app = win32com.client.GetActiveObject("PowerPoint.Application")
                app_path = str(app.Path).lower()
                if "kingsoft" not in app_path and "wps" not in app_path:
                    return app
            except Exception:
                pass
        raise RuntimeError("启动 PowerPoint 超时，请检查 POWERPNT_EXE 路径")

    # 回退：直接 Dispatch（仅 WPS 可用时）
    return win32com.client.Dispatch("PowerPoint.Application")


def assemble(plan, output_name=None):
    try:
        import win32com.client  # noqa: F401
    except ImportError:
        raise ImportError("需要 pywin32：pip install pywin32")

    if not output_name:
        output_name = datetime.now().strftime("%Y%m%d_%H%M%S") + "_assembled"

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = OUTPUT_DIR / f"{output_name}.pptx"

    pptApp = _get_ppt_app()
    pptApp.Visible = True
    pptApp.DisplayAlerts = 0  # ppAlertsNone，屏蔽保存时的"压缩图片"等弹窗

    # 自动检测颜色修复方向
    template_is_light = "浅色底" in str(TEMPLATE_PATH)
    template_is_dark  = "深色底" in str(TEMPLATE_PATH)
    fix_light_indices = []  # 深色底来源 → 浅色底模板，需修复白字为深色
    fix_dark_indices  = []  # 浅色底来源 → 深色底模板，需修复深字为白色

    pres = None
    try:
        print(f"[组装] 共 {len(plan)} 页\n")

        pres = pptApp.Presentations.Add(WithWindow=True)
        while pres.Slides.Count > 0:
            pres.Slides(1).Delete()
        pres.Windows(1).Activate()
        pres.Windows(1).ViewType = PP_VIEW_NORMAL

        for i, item in enumerate(plan, 1):
            pptApp.DisplayAlerts = 0  # 每次迭代重置，防止 Open/Close 过程中被重置
            replace_title = item.get("replace_title")

            # ══ 纯模板页（封面 / 过渡页 / 封底） ════════════════════════════
            if "template_page" in item:
                tmpl_page = item["template_page"]
                print(f"  P{i:02d}: [模板 P{tmpl_page}]" + (f"  →  「{replace_title}」" if replace_title else ""))

                tmpl_pres = pptApp.Presentations.Open(
                    str(TEMPLATE_PATH.resolve()), ReadOnly=True, Untitled=True, WithWindow=False
                )
                pptApp.DisplayAlerts = 0
                tmpl_pres.Slides(tmpl_page).Copy()
                tmpl_pres.Close()

                count_before = pres.Slides.Count
                paste_slide_with_source_format(pptApp, pres, count_before)

                if replace_title:
                    set_template_title(pres.Slides(pres.Slides.Count), replace_title)
                continue

            # ══ 整页复制（保留源格式，不套模板框架，用于保留封面/封底） ══════
            if "copy_slide" in item:
                src  = Path(item["copy_slide"]).resolve()
                page = item["page"]
                print(f"  P{i:02d}: [整页保留] {src.name}  第{page}页")

                src_pres = pptApp.Presentations.Open(
                    str(src), ReadOnly=True, Untitled=True, WithWindow=False
                )
                pptApp.DisplayAlerts = 0
                src_pres.Slides(page).Copy()
                src_pres.Close()

                count_before = pres.Slides.Count
                paste_slide_with_source_format(pptApp, pres, count_before)
                continue

            # ══ 内容页（模板背景 + 源内容形状） ══════════════════════════════
            src  = Path(item["src"]).resolve()
            page = item["page"]
            print(f"  P{i:02d}: {src.name}  第{page}页")

            # Step 1：打开源文件，检测标题 + 确定内容形状索引（不复制到剪贴板）
            src_pres = pptApp.Presentations.Open(
                str(src), ReadOnly=True, Untitled=True, WithWindow=False
            )
            pptApp.DisplayAlerts = 0
            src_slide = src_pres.Slides(page)

            # 根据源文件是否有标题栏决定使用哪个模板页（与 replace_title 无关）
            source_title, title_shape_idx = get_source_title(src_slide)
            source_has_title = source_title is not None

            if source_has_title:
                tmpl_page       = TEMPLATE_CONTENT_PAGE           # P3：带标题栏
                content_indices = get_content_indices(src_slide, exclude_idx=title_shape_idx)
                title_text      = replace_title or source_title    # 允许 replace_title 覆盖
            else:
                tmpl_page       = TEMPLATE_CONTENT_PAGE_NO_TITLE  # P4：无标题栏
                content_indices = list(range(1, src_slide.Shapes.Count + 1))  # 复制全部形状
                title_text      = None  # 无标题栏，不写标题

            src_pres.Close()  # 先关闭，不复制（避免覆盖后面的模板页剪贴板）

            # Step 2：复制对应模板页 → 粘贴（建立背景）
            tmpl_pres = pptApp.Presentations.Open(
                str(TEMPLATE_PATH.resolve()), ReadOnly=True, Untitled=True, WithWindow=False
            )
            pptApp.DisplayAlerts = 0
            tmpl_pres.Slides(tmpl_page).Copy()
            tmpl_pres.Close()

            count_before = pres.Slides.Count
            paste_slide_with_source_format(pptApp, pres, count_before)
            current_idx = pres.Slides.Count

            # Step 3：重新打开源文件，复制内容形状 → 粘贴（保留源格式/颜色）
            if content_indices:
                src_pres = pptApp.Presentations.Open(
                    str(src), ReadOnly=True, Untitled=True, WithWindow=False
                )
                pptApp.DisplayAlerts = 0
                src_pres.Slides(page).Shapes.Range(content_indices).Copy()
                src_pres.Close()
                paste_shapes_with_source_format(pptApp, pres, current_idx)
            else:
                print(f"    [信息] 未找到内容形状")

            # Step 4：写入标题（仅带标题栏的页面）
            if source_has_title:
                set_template_title(pres.Slides(current_idx), title_text)
                print(f"        [P{tmpl_page} 带标题栏]  标题: 「{title_text}」")
            else:
                print(f"        [P{tmpl_page} 无标题栏]")

            # Step 5：标记需要颜色修复的页（自动检测 + 支持显式指定）
            src_is_dark  = "深色底" in str(item.get("src", ""))
            src_is_light = "浅色底" in str(item.get("src", ""))
            if item.get("fix_colors",      template_is_light and src_is_dark):
                fix_light_indices.append(current_idx)
            if item.get("fix_colors_dark", template_is_dark  and src_is_light):
                fix_dark_indices.append(current_idx)

        # 保存（24 = ppSaveAsOpenXMLPresentation，明确格式可屏蔽压缩图片弹窗）
        pptApp.DisplayAlerts = 0
        pres.SaveAs(str(output_path.resolve()), 24)
        print(f"\n[完成] {output_path}")

    finally:
        try:
            if pres is not None:
                pres.Saved = True
                pres.Close()
        except Exception:
            pass
        try:
            pptApp.Quit()
        except Exception:
            pass

    # win32com 关闭后，用 python-pptx 修复颜色（文件已解锁）
    if fix_light_indices:
        fix_colors_for_light_template(output_path, fix_light_indices)
    if fix_dark_indices:
        fix_colors_for_dark_template(output_path, fix_dark_indices)

    return output_path


if __name__ == "__main__":
    print("assemble_template.py 是纯引擎，请通过任务脚本或 Skill 调用 assemble(plan, output_name)")
    print("示例：")
    print("  from assemble_template import assemble")
    print("  assemble(plan, '输出文件名')")

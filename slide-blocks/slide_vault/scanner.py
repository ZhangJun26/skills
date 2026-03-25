"""
scanner.py - PPT 扫描 & 文字提取模块
"""

import os
import sys
import hashlib
import sqlite3
import json
from datetime import datetime
from pathlib import Path

# Windows 终端强制 UTF-8 输出
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .config import get_db_path, get_materials_dir, load_config


# ─── 数据库初始化 ────────────────────────────────────────────────

def init_db():
    conn = sqlite3.connect(get_db_path())
    conn.execute("""
        CREATE TABLE IF NOT EXISTS slides (
            id            INTEGER PRIMARY KEY AUTOINCREMENT,
            file_path     TEXT NOT NULL,
            file_name     TEXT NOT NULL,
            file_hash     TEXT NOT NULL,
            slide_index   INTEGER NOT NULL,
            title         TEXT,
            body_text     TEXT,
            shape_count   INTEGER,
            has_image     BOOLEAN,
            has_chart     BOOLEAN,
            file_mtime    TEXT,
            indexed_at    TEXT,
            UNIQUE(file_hash, slide_index)
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS tags (
            slide_id      INTEGER PRIMARY KEY,
            scene         TEXT,
            content_type  TEXT,
            industries    TEXT,
            keywords      TEXT,
            quality_score INTEGER,
            summary       TEXT,
            tagged_at     TEXT,
            FOREIGN KEY(slide_id) REFERENCES slides(id)
        )
    """)
    conn.commit()
    conn.close()
    print(f"[DB] 数据库已初始化：{get_db_path()}")


# ─── 文件 Hash ────────────────────────────────────────────────────

def file_hash(path: Path) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        while chunk := f.read(8192):
            h.update(chunk)
    return h.hexdigest()


# ─── 单页内容提取 ────────────────────────────────────────────────

def extract_slide_content(slide) -> dict:
    texts = []
    title = None
    has_image = False
    has_chart = False
    shape_count = len(slide.shapes)

    for shape in slide.shapes:
        # 标题优先取 placeholder 类型为 title 的
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # 判断是否是标题占位符
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PLACEHOLDER
                pass
            try:
                ph = shape.placeholder_format
                if ph and ph.idx == 0:  # idx=0 通常是标题
                    title = text
                else:
                    texts.append(text)
            except Exception:
                texts.append(text)

        if shape.shape_type == 13:  # PICTURE
            has_image = True
        try:
            _ = shape.chart
            has_chart = True
        except Exception:
            pass
        if shape.shape_type == 3:  # LINKED_OLE_OBJECT / table
            pass

    # 再扫一遍图片和图表
    for shape in slide.shapes:
        if shape.shape_type == 13:
            has_image = True
        if shape.has_text_frame is False and hasattr(shape, "image"):
            has_image = True
        try:
            _ = shape.chart
            has_chart = True
        except Exception:
            pass
        if shape.shape_type == 19:  # TABLE
            has_chart = True

    # 如果 title 没有从 placeholder 里找到，退而用第一段文字
    if not title and texts:
        title = texts[0]
        texts = texts[1:]

    body_text = "\n".join(texts).replace("\xa0", " ")
    return {
        "title": title,
        "body_text": body_text,
        "shape_count": shape_count,
        "has_image": has_image,
        "has_chart": has_chart,
    }


# ─── 扫描单个文件 ────────────────────────────────────────────────

def scan_file(pptx_path: Path, conn: sqlite3.Connection = None) -> list[dict]:
    if conn is None:
        conn = sqlite3.connect(get_db_path())
    fhash = file_hash(pptx_path)
    fmtime = datetime.fromtimestamp(pptx_path.stat().st_mtime).isoformat()
    fname = pptx_path.name

    # 检查是否已完整索引（同 hash 的所有页都存在则跳过）
    existing = conn.execute(
        "SELECT COUNT(*) FROM slides WHERE file_hash = ?", (fhash,)
    ).fetchone()[0]

    try:
        prs = Presentation(str(pptx_path))
        total = len(prs.slides)
    except Exception as e:
        print(f"  [!] 无法打开 {fname}：{e}")
        return []

    if existing >= total:
        print(f"  [跳过] {fname}（已索引 {existing} 页）")
        return []

    results = []
    now = datetime.now().isoformat()

    for i, slide in enumerate(prs.slides, start=1):
        content = extract_slide_content(slide)
        try:
            cursor = conn.execute(
                """
                INSERT OR IGNORE INTO slides
                (file_path, file_name, file_hash, slide_index,
                 title, body_text, shape_count, has_image, has_chart,
                 file_mtime, indexed_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    str(pptx_path),
                    fname,
                    fhash,
                    i,
                    content["title"],
                    content["body_text"],
                    content["shape_count"],
                    content["has_image"],
                    content["has_chart"],
                    fmtime,
                    now,
                ),
            )
            if cursor.rowcount > 0:
                results.append({"slide_index": i, **content})
        except Exception as e:
            print(f"  [!] 第{i}页写入失败：{e}")

    conn.commit()
    return results


# ─── 批量扫描目录 ────────────────────────────────────────────────

def scan_directory(directory: Path = None):
    if directory is None:
        directory = get_materials_dir()

    cfg = load_config()
    exclude_dirs = set(cfg.get("exclude_dirs", []))

    init_db()
    conn = sqlite3.connect(get_db_path())

    # 收集文件，跳过排除目录
    pptx_files = []
    for f in directory.rglob("*.pptx"):
        if any(part in exclude_dirs for part in f.parts):
            continue
        pptx_files.append(f)

    print(f"\n[扫描] 发现 {len(pptx_files)} 个文件（已排除：{exclude_dirs}）\n")

    total_slides = 0
    for pptx_path in pptx_files:
        print(f"[处理] {pptx_path.name}")
        slides = scan_file(pptx_path, conn)
        print(f"  → 新增 {len(slides)} 页")
        total_slides += len(slides)

    conn.close()
    print(f"\n[完成] 共新增 {total_slides} 条 slide 记录")


# ─── 预览提取结果 ────────────────────────────────────────────────

def preview_results(limit: int = 10):
    conn = sqlite3.connect(get_db_path())
    rows = conn.execute(
        """
        SELECT slide_index, title, body_text, shape_count, has_image, has_chart
        FROM slides
        ORDER BY slide_index
        LIMIT ?
        """,
        (limit,),
    ).fetchall()
    conn.close()

    print(f"\n{'='*60}")
    print(f"前 {limit} 页提取结果预览")
    print(f"{'='*60}")
    for row in rows:
        idx, title, body, shapes, img, chart = row
        body_preview = (body or "")[:100].replace("\n", " | ")
        print(f"\n【第 {idx} 页】")
        print(f"  标题：{title or '（无）'}")
        print(f"  正文：{body_preview or '（无文字）'}")
        print(f"  元素数：{shapes}  含图片：{'是' if img else '否'}  含图表：{'是' if chart else '否'}")


if __name__ == "__main__":
    scan_directory()
    preview_results(limit=20)
